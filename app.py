from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional
from io import BytesIO
from pathlib import Path
import os

from pptx import Presentation  # python-pptx
from urllib.parse import quote
from tempfile import NamedTemporaryFile
import subprocess
import shutil


class ReportPayload(BaseModel):
    SERVICE_CODE: str
    ID_NUMBER: str
    NAME_AR: str
    NAME_EN: str
    DAYS_COUNT: int
    ENTRY_DATE_GREGORIAN: str
    EXIT_DATE_GREGORIAN: str
    ENTRY_DATE_HIJRI: Optional[str] = None
    EXIT_DATE_HIJRI: Optional[str] = None
    REPORT_ISSUE_DATE: str
    NATIONALITY_AR: str
    NATIONALITY_EN: str
    DOCTOR_NAME_AR: str
    DOCTOR_NAME_EN: str
    JOB_TITLE_AR: str
    JOB_TITLE_EN: str
    HOSPITAL_NAME_AR: str
    HOSPITAL_NAME_EN: str
    PRINT_DATE: str
    PRINT_TIME: str


def get_template_path() -> Path:
    # Allow override via env; default to repo's public/templates/report_template.pptx
    env_path = os.getenv("PPTX_TEMPLATE_PATH")
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p
    # backend/app.py -> repo_root/backend ; go to public/templates/
    repo_root = Path(__file__).resolve().parent
    default_path = repo_root / "public" / "templates" / "report_template.pptx"
    return default_path


def replace_placeholders(prs: Presentation, mapping: dict):
    # Replace text placeholders in all shapes across all slides
    # Do replacements per-run to preserve formatting (font color/size)
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text or ""
                            new_text = text
                            for key, value in mapping.items():
                                new_text = new_text.replace(f"{{{{{key}}}}}", str(value) if value is not None else "")
                            if new_text != text:
                                run.text = new_text  # preserves run formatting
            except Exception:
                # Skip shapes that fail to process to avoid taking down the request
                continue


app = FastAPI(title="PPTX Generator Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # adjust in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/generate-pptx")
def generate_pptx(payload: ReportPayload):
    template_path = get_template_path()
    if not template_path.exists():
        raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")

    try:
        prs = Presentation(str(template_path))

        # Build mapping from placeholders to values.
        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": payload.ENTRY_DATE_GREGORIAN,
            "EXIT_DATE_GREGORIAN": payload.EXIT_DATE_GREGORIAN,
            "ENTRY_DATE_HIJRI": payload.ENTRY_DATE_HIJRI,
            "EXIT_DATE_HIJRI": payload.EXIT_DATE_HIJRI,
            "REPORT_ISSUE_DATE": payload.REPORT_ISSUE_DATE,
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": payload.PRINT_DATE,
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        filename = f"sickLeaves_{payload.NAME_AR}_{payload.ID_NUMBER}.pptx"
        # HTTP headers must be latin-1 encodable in Starlette; use RFC5987 filename*
        ascii_fallback = "sickLeaves.pptx"
        content_disposition = (
            f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        )
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": content_disposition
            },
        )
    except Exception as e:
        # Log server-side for debugging
        import traceback
        print("[generate-pptx] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-pdf")
def generate_pdf(payload: ReportPayload):
    template_path = get_template_path()
    if not template_path.exists():
        raise HTTPException(status_code=500, detail=f"Template not found: {template_path}")

    try:
        prs = Presentation(str(template_path))

        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": payload.ENTRY_DATE_GREGORIAN,
            "EXIT_DATE_GREGORIAN": payload.EXIT_DATE_GREGORIAN,
            "ENTRY_DATE_HIJRI": payload.ENTRY_DATE_HIJRI,
            "EXIT_DATE_HIJRI": payload.EXIT_DATE_HIJRI,
            "REPORT_ISSUE_DATE": payload.REPORT_ISSUE_DATE,
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": payload.PRINT_DATE,
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        # Save to temp PPTX then convert to PDF
        with NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_pptx:
            prs.save(tmp_pptx.name)
            tmp_pptx_path = tmp_pptx.name

        # Convert PPTX -> PDF using LibreOffice (headless soffice)
        # We avoid requiring Microsoft PowerPoint.
        try:
            tmp_dir = Path(tmp_pptx_path).parent
            soffice_path = os.getenv("LIBREOFFICE_PATH") or shutil.which("soffice") or shutil.which("soffice.exe")
            if not soffice_path:
                raise RuntimeError("LibreOffice 'soffice' not found in PATH. Set LIBREOFFICE_PATH or install LibreOffice.")

            # Use --outdir to place the PDF next to the PPTX
            cmd = [
                soffice_path,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp_dir),
                str(tmp_pptx_path),
            ]
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr or result.stdout}")

            # Determine output path (LibreOffice names it with .pdf in same base name)
            produced_pdf = Path(tmp_pptx_path).with_suffix(".pdf")
            if not produced_pdf.exists():
                raise RuntimeError("Converted PDF not found after LibreOffice run.")

            with open(produced_pdf, "rb") as f:
                pdf_bytes = f.read()
        except Exception as conv_err:
            raise HTTPException(status_code=500, detail=str(conv_err))

        # Stream PDF
        try:
            with open(tmp_pdf_path, "rb") as f:
                pdf_bytes = f.read()
        finally:
            # Cleanup temp files
            try:
                os.remove(tmp_pptx_path)
            except Exception:
                pass
            try:
                os.remove(Path(tmp_pptx_path).with_suffix(".pdf"))
            except Exception:
                pass

        filename = f"sickLeaves_{payload.NAME_AR}_{payload.ID_NUMBER}.pdf"
        ascii_fallback = "sickLeaves.pdf"
        cd = f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        return StreamingResponse(
            BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": cd},
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print("[generate-pdf] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 8000)))



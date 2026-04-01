from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional
from io import BytesIO
from pathlib import Path
import os
import re
from dotenv import load_dotenv

# Load environment variables from .env.local
load_dotenv(dotenv_path=".env.local")

from pptx import Presentation  # python-pptx
from urllib.parse import quote
import requests


class ReportPayload(BaseModel):
    SERVICE_CODE: str
    ID_NUMBER: str
    NAME_AR: str
    NAME_EN: str
    DAYS_COUNT: int
    ENTRY_DATE_GREGORIAN: str
    EXIT_DATE_GREGORIAN: str
    ENTRY_DATE_HIJRI: Optional[str] = None
    EXIT_DATE_HIJRI: Optional[str] = None
    REPORT_ISSUE_DATE: str
    NATIONALITY_AR: str
    NATIONALITY_EN: str
    DOCTOR_NAME_AR: str
    DOCTOR_NAME_EN: str
    JOB_TITLE_AR: str
    JOB_TITLE_EN: str
    HOSPITAL_NAME_AR: str
    HOSPITAL_NAME_EN: str
    PRINT_DATE: str
    PRINT_TIME: str


def get_template_path() -> Path:
    # Allow override via env; default to locating near current file or project root
    env_path = os.getenv("PPTX_TEMPLATE_PATH")
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p

    current_dir = Path(__file__).resolve().parent
    candidates = [
        # Vercel serverless common locations first
        Path("/var/task/backend/public/templates/report_template.pptx"),
        Path("/var/task/public/templates/report_template.pptx"),
        # Typical when backend is deployed as its own project (template inside backend/public/...)
        current_dir / "public" / "templates" / "report_template.pptx",
        # Typical when backend is in a subfolder and template is at repo root public/templates
        current_dir.parent / ".." / "public" / "templates" / "report_template.pptx",
        # When includeFiles uses repo-root paths like backend/public/**
        current_dir / ".." / "public" / "templates" / "report_template.pptx",
        Path("backend/public/templates/report_template.pptx").resolve(),
        # Relative to CWD as last resort
        Path("public/templates/report_template.pptx"),
        Path("backend/public/templates/report_template.pptx"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    # As a last resort, search for the file anywhere under working dir
    try:
        for match in Path.cwd().rglob("report_template.pptx"):
            return match
    except Exception:
        pass
    # Fall back; caller will validate existence
    return candidates[0]


def format_date_dd_mm_yyyy(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    s = str(value).strip()
    # Find first occurrence of YYYY-MM-DD or YYYY/MM/DD anywhere in the string (e.g., ISO timestamps)
    m = re.search(r"(\d{4})[-/](\d{1,2})[-/](\d{1,2})", s)
    if not m:
        return s
    yyyy, mm, dd = m.groups()
    mm = mm.zfill(2)
    dd = dd.zfill(2)
    return f"{dd}-{mm}-{yyyy}"


def load_template_presentation() -> Presentation:
    """Load template as Presentation from local file or URL."""
    # Priority 1: Local file
    local_path = get_template_path()
    if local_path and local_path.exists():
        try:
            return Presentation(str(local_path))
        except Exception as e:
            print(f"Error loading local template {local_path}: {e}")

    # Priority 2: External URL fallback
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    if not template_url:
        raise HTTPException(
            status_code=500, 
            detail="Local template not found and PPTX_TEMPLATE_URL is not set"
        )
    
    try:
        resp = requests.get(template_url, timeout=20)
        if resp.status_code != 200:
            raise HTTPException(status_code=500, detail=f"Failed to fetch template from URL: {resp.status_code}")
        return Presentation(BytesIO(resp.content))
    except requests.RequestException as e:
        raise HTTPException(status_code=500, detail=f"Error fetching template from URL: {str(e)}")


def replace_placeholders(prs: Presentation, mapping: dict):
    # Replace text placeholders in all shapes across all slides
    # Do replacements per-run to preserve formatting (font color/size)
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text or ""
                            new_text = text
                            for key, value in mapping.items():
                                new_text = new_text.replace(f"{{{{{key}}}}}", str(value) if value is not None else "")
                            if new_text != text:
                                run.text = new_text  # preserves run formatting
            except Exception:
                # Skip shapes that fail to process to avoid taking down the request
                continue


app = FastAPI(title="PPTX Generator Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # adjust in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/")
@app.get("/health")
def health():
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    return {
        "status": "ok",
        "template_url_configured": bool(template_url),
        "template_url": template_url if template_url else "Not configured"
    }


@app.get("/debug-template")
def debug_template():
    p = get_template_path()
    candidates = []
    try:
        roots = [Path("/var/task"), Path.cwd()]
        found = []
        for root in roots:
            try:
                for match in root.rglob("report_template.pptx"):
                    found.append(str(match))
            except Exception:
                continue
    except Exception:
        found = []
    checks = {
        "/var/task/backend/public/templates": (Path("/var/task/backend/public/templates").exists()),
        "/var/task/public/templates": (Path("/var/task/public/templates").exists()),
    }
    template_url = os.getenv("PPTX_TEMPLATE_URL")
    return {
        "resolved_path": str(p),
        "exists": p.exists(),
        "cwd": str(Path.cwd()),
        "file_dir": str(Path(__file__).resolve().parent),
        "found_candidates": found,
        "dir_checks": checks,
        "template_url": template_url or None,
    }


@app.post("/generate-pptx")
def generate_pptx(payload: ReportPayload):
    try:
        prs = load_template_presentation()

        # Build mapping from placeholders to values.
        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.ENTRY_DATE_GREGORIAN),
            "EXIT_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.EXIT_DATE_GREGORIAN),
            "ENTRY_DATE_HIJRI": format_date_dd_mm_yyyy(payload.ENTRY_DATE_HIJRI),
            "EXIT_DATE_HIJRI": format_date_dd_mm_yyyy(payload.EXIT_DATE_HIJRI),
            "REPORT_ISSUE_DATE": format_date_dd_mm_yyyy(payload.REPORT_ISSUE_DATE),
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": format_date_dd_mm_yyyy(payload.PRINT_DATE),
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)

        filename = "sickLeaves.pptx"
        # HTTP headers must be latin-1 encodable in Starlette; use RFC5987 filename*
        ascii_fallback = "sickLeaves.pptx"
        content_disposition = (
            f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        )
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": content_disposition
            },
        )
    except Exception as e:
        # Log server-side for debugging
        import traceback
        print("[generate-pptx] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


SLIDIZE_CONVERT_URL = "https://api.slidize.cloud/v1.0/slides/convert/pdf"


@app.post("/generate-pdf")
def generate_pdf(payload: ReportPayload):
    try:
        prs = load_template_presentation()

        mapping = {
            "SERVICE_CODE": payload.SERVICE_CODE,
            "ID_NUMBER": payload.ID_NUMBER,
            "NAME_AR": payload.NAME_AR,
            "NAME_EN": payload.NAME_EN,
            "DAYS_COUNT": payload.DAYS_COUNT,
            "ENTRY_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.ENTRY_DATE_GREGORIAN),
            "EXIT_DATE_GREGORIAN": format_date_dd_mm_yyyy(payload.EXIT_DATE_GREGORIAN),
            "ENTRY_DATE_HIJRI": format_date_dd_mm_yyyy(payload.ENTRY_DATE_HIJRI),
            "EXIT_DATE_HIJRI": format_date_dd_mm_yyyy(payload.EXIT_DATE_HIJRI),
            "REPORT_ISSUE_DATE": format_date_dd_mm_yyyy(payload.REPORT_ISSUE_DATE),
            "NATIONALITY_AR": payload.NATIONALITY_AR,
            "NATIONALITY_EN": payload.NATIONALITY_EN,
            "DOCTOR_NAME_AR": payload.DOCTOR_NAME_AR,
            "DOCTOR_NAME_EN": payload.DOCTOR_NAME_EN,
            "JOB_TITLE_AR": payload.JOB_TITLE_AR,
            "JOB_TITLE_EN": payload.JOB_TITLE_EN,
            "HOSPITAL_NAME_AR": payload.HOSPITAL_NAME_AR,
            "HOSPITAL_NAME_EN": payload.HOSPITAL_NAME_EN,
            "PRINT_DATE": format_date_dd_mm_yyyy(payload.PRINT_DATE),
            "PRINT_TIME": payload.PRINT_TIME,
        }

        replace_placeholders(prs, mapping)

        # Save filled PPTX to in-memory buffer
        pptx_buf = BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)

        # Convert PPTX -> PDF using Slidize Cloud API (free, no auth required)
        try:
            slidize_resp = requests.post(
                SLIDIZE_CONVERT_URL,
                files={"documents": ("report.pptx", pptx_buf, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                timeout=120,
            )
            if slidize_resp.status_code != 200:
                raise RuntimeError(
                    f"Slidize Cloud API returned status {slidize_resp.status_code}: {slidize_resp.text[:500]}"
                )
            pdf_bytes = slidize_resp.content
        except requests.RequestException as e:
            raise HTTPException(status_code=500, detail=f"Error calling Slidize Cloud API: {str(e)}")

        filename = "sickLeaves.pdf"
        ascii_fallback = "sickLeaves.pdf"
        cd = f"attachment; filename=\"{ascii_fallback}\"; filename*=UTF-8''{quote(filename)}"
        return StreamingResponse(
            BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": cd},
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print("[generate-pdf] Error:", e)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
